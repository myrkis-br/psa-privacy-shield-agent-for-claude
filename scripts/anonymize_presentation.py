"""
PSA - Privacy Shield Agent
Script: anonymize_presentation.py
Responsável: PSA Guardião

Anonimiza apresentações PPTX:
  - Extrai texto de cada slide (título + corpo + notas)
  - Detecta e substitui CPF, CNPJ, email, telefone, CEP, datas, valores monetários
  - Substitui nomes por tokens [PESSOA_N] / [EMPRESA_N]
  - Salva como JSON estruturado em data/anonymized/
  - Salva mapa de entidades em data/maps/

Formato de saída JSON:
  {
    "metadados": { ... },
    "slides": {
      "slide_1": { "titulo": "...", "corpo": ["...", "..."], "notas": "..." },
      ...
    }
  }

Uso:
  python3 scripts/anonymize_presentation.py <caminho_do_arquivo> [--slides <n>]

Exemplos:
  python3 scripts/anonymize_presentation.py data/real/pitch.pptx
  python3 scripts/anonymize_presentation.py data/real/relatorio.pptx --slides 10
"""

import sys
import json
import argparse
import logging
from pathlib import Path
from datetime import datetime
from typing import List, Dict, Optional, Tuple

# ---------------------------------------------------------------------------
# Configuração de caminhos
# ---------------------------------------------------------------------------

BASE_DIR = Path(__file__).resolve().parent.parent
ANONYMIZED_DIR = BASE_DIR / "data" / "anonymized"
MAPS_DIR = BASE_DIR / "data" / "maps"
LOGS_DIR = BASE_DIR / "logs"

for d in (ANONYMIZED_DIR, MAPS_DIR, LOGS_DIR):
    d.mkdir(parents=True, exist_ok=True)

sys.path.insert(0, str(BASE_DIR / "scripts"))

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [PSA-Guardião] %(levelname)s: %(message)s",
    handlers=[
        logging.FileHandler(LOGS_DIR / "psa_guardiao.log"),
        logging.StreamHandler(sys.stdout),
    ],
)
log = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Imports de terceiros
# ---------------------------------------------------------------------------

try:
    from pptx import Presentation
    from pptx.util import Pt
except ImportError:
    Presentation = None

from text_engine import TextAnonymizer

# ---------------------------------------------------------------------------
# Extração de conteúdo por slide
# ---------------------------------------------------------------------------

def _get_shape_text(shape) -> Optional[str]:
    """Extrai texto de uma shape (caixa de texto, tabela, etc.)."""
    try:
        if shape.has_text_frame:
            lines = []
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    lines.append(text)
            return "\n".join(lines) if lines else None

        # Tabelas dentro de slides
        if shape.has_table:
            rows = []
            for row in shape.table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    rows.append(" | ".join(cells))
            return "\n".join(rows) if rows else None

    except Exception:
        pass
    return None


def _extract_slide(slide, slide_num: int) -> Dict:
    """
    Extrai todo o conteúdo de um slide.
    Retorna dict com titulo, corpo (lista de textos), notas.
    """
    titulo = None
    corpo_items = []
    notas = None

    # Tenta extrair título pelo placeholder
    try:
        if slide.shapes.title and slide.shapes.title.text.strip():
            titulo = slide.shapes.title.text.strip()
    except Exception:
        pass

    # Extrai demais shapes
    for shape in slide.shapes:
        # Pula o título (já extraído)
        try:
            if shape == slide.shapes.title:
                continue
        except Exception:
            pass

        text = _get_shape_text(shape)
        if text:
            corpo_items.append(text)

    # Notas do apresentador
    try:
        notes_slide = slide.notes_slide
        if notes_slide:
            notes_tf = notes_slide.notes_text_frame
            if notes_tf:
                notas_text = notes_tf.text.strip()
                if notas_text:
                    notas = notas_text
    except Exception:
        pass

    return {
        "slide_num": slide_num,
        "titulo": titulo or "",
        "corpo": corpo_items,
        "notas": notas or "",
        "word_count": len((titulo or "").split()) + sum(len(t.split()) for t in corpo_items),
    }


def _extract_all_slides(path: Path) -> Tuple[List[Dict], Dict]:
    """
    Extrai conteúdo de todos os slides do PPTX.
    Retorna (slides, metadados_apresentacao).
    """
    if Presentation is None:
        raise ImportError("python-pptx não instalado. Execute: pip3 install python-pptx")

    prs = Presentation(str(path))
    slides = []

    # Metadados da apresentação
    core_props = prs.core_properties
    metadata = {
        "total_slides": len(prs.slides),
        "autor": core_props.author or "",
        "titulo_apresentacao": core_props.title or "",
        "assunto": core_props.subject or "",
        "data_criacao": str(core_props.created) if core_props.created else "",
        "ultima_modificacao": str(core_props.modified) if core_props.modified else "",
    }

    for i, slide in enumerate(prs.slides, start=1):
        slide_data = _extract_slide(slide, i)
        slides.append(slide_data)
        log.info(
            f"  Slide {i:02d}: título='{slide_data['titulo'][:40]}' | "
            f"{len(slide_data['corpo'])} item(ns) | {slide_data['word_count']} palavras"
        )

    return slides, metadata


def _sample_slides(slides: List[Dict], max_slides: int) -> List[Dict]:
    """
    Seleciona amostra de slides.

    Estratégia:
    - Slide 1 sempre incluído (capa / título)
    - Slide com mais conteúdo de cada seção
    - Distribui ao longo da apresentação
    """
    if len(slides) <= max_slides:
        return slides

    # Slide 1 sempre entra
    must_include = [slides[0]]
    rest = slides[1:]

    remaining = max_slides - 1
    if remaining <= 0:
        return must_include

    # Distribui uniformemente
    step = max(1, len(rest) // remaining)
    sampled = rest[::step][:remaining]

    combined = must_include + sampled
    combined.sort(key=lambda s: s["slide_num"])
    return combined[:max_slides]


# ---------------------------------------------------------------------------
# Função principal
# ---------------------------------------------------------------------------

def anonymize_presentation(
    input_path: Path,
    max_slides: int = 15,
) -> Tuple[Path, Path]:
    """
    Anonimiza uma apresentação PPTX.

    Args:
        input_path: Caminho para o PPTX em data/real/
        max_slides: Número máximo de slides na amostra

    Returns:
        Tuple (caminho_anonimizado, caminho_mapa)
    """
    input_path = Path(input_path)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    stem = input_path.stem

    log.info(f"Iniciando anonimização de apresentação: {input_path.name}")

    # --- Extração ---
    all_slides, presentation_meta = _extract_all_slides(input_path)
    total_slides = len(all_slides)
    log.info(f"Apresentação carregada: {total_slides} slides")

    # Anonimiza metadados da própria apresentação
    engine = TextAnonymizer()
    if presentation_meta.get("autor"):
        presentation_meta["autor"] = engine.anonymize(presentation_meta["autor"])
    if presentation_meta.get("titulo_apresentacao"):
        presentation_meta["titulo_apresentacao"] = engine.anonymize(
            presentation_meta["titulo_apresentacao"]
        )

    # --- Amostragem ---
    sample = _sample_slides(all_slides, max_slides)
    sampled_nums = [s["slide_num"] for s in sample]
    log.info(f"Amostra: slides {sampled_nums} ({len(sample)} de {total_slides})")

    # --- Anonimização ---
    anonymized_slides = {}

    for slide in sample:
        titulo_anon = engine.anonymize(slide["titulo"]) if slide["titulo"] else ""

        corpo_anon = []
        for item in slide["corpo"]:
            corpo_anon.append(engine.anonymize(item))

        notas_anon = engine.anonymize(slide["notas"]) if slide["notas"] else ""

        key = f"slide_{slide['slide_num']}"
        anonymized_slides[key] = {
            "titulo": titulo_anon,
            "corpo": corpo_anon,
            "notas": notas_anon,
        }

    log.info(f"Entidades substituídas: {len(engine.entity_map)} ocorrências únicas")
    for token, original in list(engine.entity_map.items())[:5]:
        log.info(f"  Exemplo: '{str(original)[:40]}' → '{token}'")

    # --- Montar JSON de saída ---
    output_data = {
        "metadados": {
            **presentation_meta,
            "gerado_em": datetime.now().isoformat(),
            "arquivo_original": input_path.name,
            "slides_incluidos": sampled_nums,
            "slides_omitidos": [
                s["slide_num"] for s in all_slides
                if s["slide_num"] not in sampled_nums
            ],
        },
        "slides": anonymized_slides,
    }

    # --- Salvar arquivo anonimizado ---
    anon_filename = f"anon_{stem}_{timestamp}.json"
    anon_path = ANONYMIZED_DIR / anon_filename
    with open(anon_path, "w", encoding="utf-8") as f:
        json.dump(output_data, f, ensure_ascii=False, indent=2)

    log.info(f"Apresentação anonimizada salva: {anon_path}")

    # --- Salvar mapa ---
    map_data = {
        "timestamp": datetime.now().isoformat(),
        "tipo": "apresentacao",
        "arquivo_original": str(input_path),
        "arquivo_anonimizado": str(anon_path),
        "total_slides_original": total_slides,
        "slides_na_amostra": sampled_nums,
        "entidades": {
            token: original
            for token, original in engine.entity_map.items()
        },
        "estatisticas": {
            "pessoas_substituidas": engine._counters["pessoa"],
            "empresas_substituidas": engine._counters["empresa"],
            "total_entidades": len(engine.entity_map),
        },
    }

    map_filename = f"map_{stem}_{timestamp}.json"
    map_path = MAPS_DIR / map_filename
    with open(map_path, "w", encoding="utf-8") as f:
        json.dump(map_data, f, ensure_ascii=False, indent=2)

    log.info(f"Mapa salvo: {map_path}")
    log.info(
        f"Anonimização concluída: {engine._counters['pessoa']} pessoas, "
        f"{engine._counters['empresa']} empresas substituídas."
    )

    return anon_path, map_path


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description="PSA Guardião — Anonimizador de apresentações PPTX"
    )
    parser.add_argument("arquivo", help="Caminho para o PPTX a anonimizar")
    parser.add_argument(
        "--slides",
        type=int,
        default=15,
        metavar="N",
        help="Número máximo de slides na amostra (padrão: 15)",
    )
    args = parser.parse_args()

    input_path = Path(args.arquivo)
    if not input_path.exists():
        log.error(f"Arquivo não encontrado: {input_path}")
        sys.exit(1)

    if input_path.suffix.lower() != ".pptx":
        log.error(f"Arquivo deve ser PPTX. Recebido: {input_path.suffix}")
        sys.exit(1)

    anon_path, map_path = anonymize_presentation(input_path, max_slides=args.slides)

    print("\n" + "=" * 60)
    print("PSA GUARDIÃO — ANONIMIZAÇÃO CONCLUÍDA")
    print("=" * 60)
    print(f"  Arquivo anonimizado    : {anon_path}")
    print(f"  Mapa de correspondência: {map_path}")
    print("=" * 60)


if __name__ == "__main__":
    main()
