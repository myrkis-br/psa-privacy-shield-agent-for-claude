"""
PSA - Privacy Shield Agent
Script: test_documents.py

Testa os 3 scripts de anonimização:
  1. anonymize_document.py → DOCX e TXT
  2. anonymize_pdf.py → PDF
  3. anonymize_presentation.py → PPTX

Para cada tipo:
  - Gera um arquivo fake com dados sensíveis em data/real/
  - Roda o anonimizador
  - Valida: entidades substituídas, mapa gerado, nenhum dado real exposto

Uso:
  python3 scripts/test_documents.py
"""

import sys
import json
import re
from pathlib import Path
from datetime import date

BASE_DIR = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(BASE_DIR / "scripts"))

REAL_DIR = BASE_DIR / "data" / "real"
REAL_DIR.mkdir(parents=True, exist_ok=True)

from faker import Faker
fake = Faker("pt_BR")
Faker.seed(1)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

RESULTADO_GLOBAL = {"passou": 0, "falhou": 0}


def check(condition: bool, msg_ok: str, msg_fail: str):
    if condition:
        print(f"   ✓ {msg_ok}")
        RESULTADO_GLOBAL["passou"] += 1
    else:
        print(f"   ✗ {msg_fail}")
        RESULTADO_GLOBAL["falhou"] += 1


def _fake_cpf() -> str:
    digits = [fake.random_int(0, 9) for _ in range(9)]
    for _ in range(2):
        s = sum((len(digits) + 1 - i) * v for i, v in enumerate(digits))
        d = (s * 10 % 11) % 10
        digits.append(d)
    return "{}{}{}.{}{}{}.{}{}{}-{}{}".format(*digits)


def _sensitive_text_block() -> str:
    """Gera bloco de texto com vários dados sensíveis embutidos."""
    nome = fake.name()
    cpf = _fake_cpf()
    email = fake.email()
    phone = fake.phone_number()
    addr = fake.street_address()
    city = fake.city()
    cep = fake.postcode()
    valor = f"R$ {fake.random_number(5):,}".replace(",", ".")
    data = date.today().strftime("%d/%m/%Y")

    return (
        f"O cliente {nome}, inscrito no CPF {cpf}, residente na {addr}, "
        f"{city}, CEP {cep}, entrou em contato pelo email {email} "
        f"e telefone {phone}. "
        f"O valor total da operação é de {valor}, realizada em {data}. "
        f"Conforme acordado com {nome}, o pagamento será efetuado em parcelas."
    )


# ---------------------------------------------------------------------------
# Teste 1: DOCX
# ---------------------------------------------------------------------------

def test_docx():
    print("\n" + "=" * 60)
    print("TESTE 1: anonymize_document.py — DOCX")
    print("=" * 60)

    try:
        from docx import Document
    except ImportError:
        print("   ⚠ python-docx não disponível, pulando teste DOCX")
        return

    from anonymize_document import anonymize_document

    # Gera DOCX fake
    doc = Document()
    doc.add_heading("Relatório de Clientes — CONFIDENCIAL", level=1)
    doc.add_paragraph(
        "Este documento contém informações sensíveis e deve ser tratado com sigilo."
    )
    doc.add_heading("1. Dados do Cliente Principal", level=2)
    doc.add_paragraph(_sensitive_text_block())
    doc.add_paragraph(_sensitive_text_block())
    doc.add_heading("2. Histórico de Transações", level=2)
    for _ in range(5):
        doc.add_paragraph(_sensitive_text_block())
    doc.add_heading("3. Observações Finais", level=2)
    doc.add_paragraph(
        f"Documento elaborado por {fake.name()}, em {date.today().strftime('%d/%m/%Y')}."
    )

    docx_path = REAL_DIR / "teste_relatorio.docx"
    doc.save(str(docx_path))
    print(f"   DOCX gerado: {docx_path}")

    # Anonimiza
    anon_path, map_path = anonymize_document(docx_path, sample_paragraphs=15)

    # Valida
    anon_text = anon_path.read_text(encoding="utf-8")
    with open(map_path) as f:
        mapa = json.load(f)

    # Lê o DOCX original para extrair os valores sensíveis gerados
    orig_doc = Document(str(docx_path))
    orig_text = "\n".join(p.text for p in orig_doc.paragraphs)
    cpf_pattern = re.compile(r'\d{3}\.\d{3}\.\d{3}-\d{2}')
    email_pattern = re.compile(r'[a-zA-Z0-9._%+\-]+@[a-zA-Z0-9.\-]+\.[a-zA-Z]{2,}')

    orig_cpfs = set(cpf_pattern.findall(orig_text))
    orig_emails = set(email_pattern.findall(orig_text))
    leaked_cpfs = [v for v in orig_cpfs if v in anon_text]
    leaked_emails = [v for v in orig_emails if v in anon_text]

    check(anon_path.exists(), "Arquivo anonimizado criado", "Arquivo não criado")
    check(map_path.exists(), "Mapa criado", "Mapa não criado")
    check(len(leaked_cpfs) == 0, "Nenhum CPF original vazou", f"CPFs originais detectados: {leaked_cpfs[:2]}")
    check(len(leaked_emails) == 0, "Nenhum email original vazou", f"Emails originais detectados: {leaked_emails[:2]}")
    check(len(mapa.get("entidades", {})) > 0, f"{len(mapa['entidades'])} entidades no mapa", "Mapa de entidades vazio")
    check("tipo" in mapa and mapa["tipo"] == "documento", "Tipo correto no mapa", "Tipo incorreto no mapa")
    check("[PESSOA_" in anon_text or "[EMPRESA_" in anon_text, "Tokens de entidade presentes", "Nenhum token de entidade encontrado")

    print(f"   → Entidades substituídas: {len(mapa['entidades'])}")
    print(f"   → Parágrafos: {mapa['total_paragrafos_amostra']} de {mapa['total_paragrafos_original']}")


# ---------------------------------------------------------------------------
# Teste 2: TXT
# ---------------------------------------------------------------------------

def test_txt():
    print("\n" + "=" * 60)
    print("TESTE 2: anonymize_document.py — TXT")
    print("=" * 60)

    from anonymize_document import anonymize_document

    # Gera TXT fake
    paragraphs = []
    paragraphs.append("CONTRATO DE PRESTAÇÃO DE SERVIÇOS\n")
    paragraphs.append(
        f"Pelo presente instrumento, {fake.name()}, CPF {_fake_cpf()}, "
        f"doravante denominado CONTRATANTE, residente na {fake.street_address()}, "
        f"cidade de {fake.city()}, CEP {fake.postcode()}, email {fake.email()}, "
        f"telefone {fake.phone_number()},\n"
    )
    paragraphs.append(
        f"CONTRATA os serviços de {fake.name()}, CPF {_fake_cpf()}, "
        f"CNPJ {fake.cnpj()}, pelo valor de R$ {fake.random_number(4):,} reais mensais,\n"
    )
    for i in range(8):
        paragraphs.append(f"Cláusula {i+1}: {_sensitive_text_block()}\n")
    paragraphs.append(
        f"Assinado em {date.today().strftime('%d de %B de %Y')}.\n"
    )

    txt_path = REAL_DIR / "teste_contrato.txt"
    txt_path.write_text("\n\n".join(paragraphs), encoding="utf-8")
    print(f"   TXT gerado: {txt_path}")

    # Anonimiza
    anon_path, map_path = anonymize_document(txt_path, sample_paragraphs=10)

    # Valida
    anon_text = anon_path.read_text(encoding="utf-8")
    with open(map_path) as f:
        mapa = json.load(f)

    cpf_pattern = re.compile(r'\d{3}\.\d{3}\.\d{3}-\d{2}')
    email_pattern = re.compile(r'[a-zA-Z0-9._%+\-]+@[a-zA-Z0-9.\-]+\.[a-zA-Z]{2,}')

    orig_text_txt = txt_path.read_text(encoding="utf-8")
    orig_cpfs = set(cpf_pattern.findall(orig_text_txt))
    orig_emails = set(email_pattern.findall(orig_text_txt))
    leaked_cpfs = [v for v in orig_cpfs if v in anon_text]
    leaked_emails = [v for v in orig_emails if v in anon_text]

    check(anon_path.exists(), "Arquivo anonimizado criado", "Arquivo não criado")
    check(len(leaked_cpfs) == 0, "Nenhum CPF original vazou", f"CPFs originais: {leaked_cpfs[:2]}")
    check(len(leaked_emails) == 0, "Nenhum email original vazou", f"Emails originais: {leaked_emails[:2]}")
    check(len(mapa.get("entidades", {})) > 0, f"{len(mapa['entidades'])} entidades no mapa", "Mapa vazio")

    print(f"   → Entidades substituídas: {len(mapa['entidades'])}")


# ---------------------------------------------------------------------------
# Teste 3: PDF
# ---------------------------------------------------------------------------

def test_pdf():
    print("\n" + "=" * 60)
    print("TESTE 3: anonymize_pdf.py — PDF")
    print("=" * 60)

    # Tenta criar PDF com reportlab (pode não estar instalado)
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.pdfgen import canvas
        has_reportlab = True
    except ImportError:
        has_reportlab = False

    if not has_reportlab:
        # Tenta criar PDF simples com fpdf2
        try:
            from fpdf import FPDF
            has_fpdf = True
        except ImportError:
            has_fpdf = False

        if not has_fpdf:
            print("   ⚠ reportlab e fpdf2 não instalados.")
            print("   Criando PDF mínimo via bytes para testar extração...")
            _test_pdf_minimal()
            return

    from anonymize_pdf import anonymize_pdf

    if has_reportlab:
        _create_pdf_reportlab(REAL_DIR / "teste_documento.pdf")
    else:
        _create_pdf_fpdf(REAL_DIR / "teste_documento.pdf")

    pdf_path = REAL_DIR / "teste_documento.pdf"
    print(f"   PDF gerado: {pdf_path}")

    anon_path, map_path = anonymize_pdf(pdf_path, max_pages=3)

    anon_text = anon_path.read_text(encoding="utf-8")
    with open(map_path) as f:
        mapa = json.load(f)

    check(anon_path.exists(), "Arquivo anonimizado criado", "Arquivo não criado")
    check(map_path.exists(), "Mapa criado", "Mapa não criado")
    check("tipo" in mapa and mapa["tipo"] == "pdf", "Tipo correto no mapa", "Tipo incorreto")
    check(len(mapa.get("paginas_na_amostra", [])) > 0, "Páginas registradas no mapa", "Nenhuma página no mapa")

    print(f"   → Páginas na amostra: {mapa['paginas_na_amostra']}")
    print(f"   → Entidades: {len(mapa.get('entidades', {}))}")


def _test_pdf_minimal():
    """Teste com PDF já existente ou pula graciosamente."""
    from anonymize_pdf import anonymize_pdf

    # Busca qualquer PDF em data/real/
    existing = list((BASE_DIR / "data" / "real").glob("*.pdf"))
    if existing:
        pdf_path = existing[0]
        print(f"   Usando PDF existente: {pdf_path}")
        anon_path, map_path = anonymize_pdf(pdf_path, max_pages=2)
        check(anon_path.exists(), "Arquivo anonimizado criado", "Arquivo não criado")
    else:
        print("   ⚠ Nenhum PDF disponível para teste. Pulando.")


def _create_pdf_reportlab(path: Path):
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas as rl_canvas
    from reportlab.lib.utils import simpleSplit

    c = rl_canvas.Canvas(str(path), pagesize=A4)
    width, height = A4

    def add_page(title: str, content: str):
        c.setFont("Helvetica-Bold", 14)
        c.drawString(50, height - 60, title)
        c.setFont("Helvetica", 11)
        y = height - 100
        for line in simpleSplit(content, "Helvetica", 11, width - 100):
            if y < 60:
                c.showPage()
                y = height - 60
                c.setFont("Helvetica", 11)
            c.drawString(50, y, line)
            y -= 16
        c.showPage()

    add_page("Relatório Confidencial — Página 1", _sensitive_text_block() * 3)
    add_page("Dados Financeiros — Página 2", _sensitive_text_block() * 3)
    add_page("Conclusões — Página 3", _sensitive_text_block() * 2)
    c.save()


def _create_pdf_fpdf(path: Path):
    from fpdf import FPDF
    pdf = FPDF()
    for i in range(3):
        pdf.add_page()
        pdf.set_font("Helvetica", "B", 14)
        pdf.cell(0, 10, f"Página {i+1}", ln=True)
        pdf.set_font("Helvetica", size=10)
        text = _sensitive_text_block()
        pdf.multi_cell(0, 8, text)
    pdf.output(str(path))


# ---------------------------------------------------------------------------
# Teste 4: PPTX
# ---------------------------------------------------------------------------

def test_pptx():
    print("\n" + "=" * 60)
    print("TESTE 4: anonymize_presentation.py — PPTX")
    print("=" * 60)

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        print("   ⚠ python-pptx não disponível, pulando teste PPTX")
        return

    from anonymize_presentation import anonymize_presentation

    # Gera PPTX fake
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # título + conteúdo

    sections = [
        ("Relatório Executivo Q1 — CONFIDENCIAL", _sensitive_text_block()),
        ("Dados de Clientes", _sensitive_text_block()),
        ("Análise Financeira", _sensitive_text_block()),
        ("Projeções", _sensitive_text_block()),
        ("Conclusões", _sensitive_text_block()),
    ]

    for title_text, body_text in sections:
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title_text
        tf = slide.placeholders[1].text_frame
        tf.text = body_text

    pptx_path = REAL_DIR / "teste_apresentacao.pptx"
    prs.save(str(pptx_path))
    print(f"   PPTX gerado: {pptx_path} ({len(sections)} slides)")

    # Anonimiza
    anon_path, map_path = anonymize_presentation(pptx_path, max_slides=10)

    # Valida
    with open(anon_path, encoding="utf-8") as f:
        anon_data = json.load(f)
    with open(map_path) as f:
        mapa = json.load(f)

    cpf_pattern = re.compile(r'\d{3}\.\d{3}\.\d{3}-\d{2}')
    email_pattern = re.compile(r'[a-zA-Z0-9._%+\-]+@[a-zA-Z0-9.\-]+\.[a-zA-Z]{2,}')

    # Lê o PPTX original para extrair os valores sensíveis gerados
    orig_prs = Presentation(str(pptx_path))
    orig_text_pptx = ""
    for sl in orig_prs.slides:
        for shape in sl.shapes:
            if shape.has_text_frame:
                orig_text_pptx += shape.text_frame.text + " "

    orig_cpfs = set(cpf_pattern.findall(orig_text_pptx))
    orig_emails = set(email_pattern.findall(orig_text_pptx))

    all_anon_text = json.dumps(anon_data, ensure_ascii=False)
    leaked_cpfs = [v for v in orig_cpfs if v in all_anon_text]
    leaked_emails = [v for v in orig_emails if v in all_anon_text]

    check(anon_path.exists(), "Arquivo JSON criado", "Arquivo não criado")
    check(map_path.exists(), "Mapa criado", "Mapa não criado")
    check("slides" in anon_data, "Estrutura 'slides' no JSON", "Estrutura 'slides' ausente")
    check("metadados" in anon_data, "Estrutura 'metadados' no JSON", "Estrutura 'metadados' ausente")
    check(len(anon_data["slides"]) > 0, f"{len(anon_data['slides'])} slides no JSON", "Nenhum slide no JSON")

    check(len(leaked_cpfs) == 0, "Nenhum CPF original vazou", f"CPFs originais: {leaked_cpfs[:2]}")
    check(len(leaked_emails) == 0, "Nenhum email original vazou", f"Emails originais: {leaked_emails[:2]}")

    check(mapa["tipo"] == "apresentacao", "Tipo correto no mapa", "Tipo incorreto")
    check(len(mapa.get("entidades", {})) > 0, f"{len(mapa['entidades'])} entidades no mapa", "Mapa vazio")

    # Verifica estrutura de cada slide
    for key, slide in anon_data["slides"].items():
        has_keys = "titulo" in slide and "corpo" in slide and "notas" in slide
        if not has_keys:
            check(False, "", f"Slide {key} sem estrutura titulo/corpo/notas")
            break
    else:
        check(True, "Todos os slides têm estrutura titulo/corpo/notas", "")

    print(f"   → Slides anonimizados: {list(anon_data['slides'].keys())}")
    print(f"   → Entidades: {len(mapa['entidades'])}")

    # Mostra amostra do JSON
    first_key = list(anon_data["slides"].keys())[0]
    first_slide = anon_data["slides"][first_key]
    print(f"\n   Amostra ({first_key}):")
    print(f"     Título: {first_slide['titulo'][:60]}")
    if first_slide["corpo"]:
        print(f"     Corpo[0]: {first_slide['corpo'][0][:80]}")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    print("=" * 60)
    print("PSA — TESTE DOS SCRIPTS DE ANONIMIZAÇÃO")
    print("=" * 60)

    test_docx()
    test_txt()
    test_pdf()
    test_pptx()

    print("\n" + "=" * 60)
    passou = RESULTADO_GLOBAL["passou"]
    falhou = RESULTADO_GLOBAL["falhou"]
    total = passou + falhou
    print(f"RESULTADO FINAL: {passou}/{total} verificações passaram")
    if falhou == 0:
        print("TODOS OS TESTES PASSARAM ✓")
    else:
        print(f"FALHAS: {falhou} verificação(ões)")
    print("=" * 60)

    sys.exit(0 if falhou == 0 else 1)


if __name__ == "__main__":
    main()
