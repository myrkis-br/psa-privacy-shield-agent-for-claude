"""
Script auxiliar: cria apresentação PPTX realista sobre a Operação Lava Jato.
Todos os nomes são FICTÍCIOS (gerados por Faker) mas cargos, datas, estrutura
e valores são baseados em informações públicas do caso.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from faker import Faker
import random
import os

fake = Faker("pt_BR")
Faker.seed(int.from_bytes(os.urandom(8), "big"))

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG_DARK = RGBColor(0x1a, 0x1a, 0x2e)
ACCENT = RGBColor(0x00, 0xe5, 0xff)
WHITE = RGBColor(0xff, 0xff, 0xff)
GRAY = RGBColor(0x99, 0x99, 0x99)
RED = RGBColor(0xff, 0x44, 0x44)
GREEN = RGBColor(0x44, 0xff, 0x88)

def add_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK

def add_title_box(slide, text, top=Inches(0.5), size=Pt(36)):
    txBox = slide.shapes.add_textbox(Inches(0.8), top, Inches(11.5), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.LEFT
    return tf

def add_body(slide, lines, top=Inches(1.8), size=Pt(16)):
    txBox = slide.shapes.add_textbox(Inches(0.8), top, Inches(11.5), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = size
        p.font.color.rgb = WHITE
        p.space_after = Pt(6)
    return tf

# --- Gerar nomes fictícios consistentes ---
nome_expresidente = fake.name_male()
nome_exministro1 = fake.name_male()
nome_exministro2 = fake.name_male()
nome_dir_petrobras1 = fake.name_male()
nome_dir_petrobras2 = fake.name_male()
nome_dir_petrobras3 = fake.name_female()
nome_doleiro1 = fake.name_male()
nome_doleiro2 = fake.name_male()
nome_empreiteiro1 = fake.name_male()
nome_empreiteiro2 = fake.name_male()
nome_empreiteiro3 = fake.name_male()
nome_delator1 = fake.name_male()
nome_delator2 = fake.name_male()
nome_delator3 = fake.name_male()
nome_juiz = fake.name_male()
nome_procurador = fake.name_male()

empresa1 = fake.company()
empresa2 = fake.company()
empresa3 = fake.company()
empresa4 = fake.company()

def fake_cnpj():
    def calc(digits, weights):
        s = sum(d * w for d, w in zip(digits, weights))
        r = s % 11
        return 0 if r < 2 else 11 - r
    base = [random.randint(0, 9) for _ in range(8)] + [0, 0, 0, 1]
    d1 = calc(base, [5, 4, 3, 2, 9, 8, 7, 6, 5, 4, 3, 2])
    d2 = calc(base + [d1], [6, 5, 4, 3, 2, 9, 8, 7, 6, 5, 4, 3, 2])
    n = base + [d1, d2]
    return "{}{}.{}{}{}.{}{}{}/{}{}{}{}-{}{}".format(*n)

cnpj1 = fake_cnpj()
cnpj2 = fake_cnpj()
cnpj3 = fake_cnpj()
cnpj4 = fake_cnpj()

cpf_expresidente = fake.cpf()
cpf_doleiro1 = fake.cpf()
cpf_delator1 = fake.cpf()

# ============================================================
# SLIDE 1 — Título
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide1)
add_title_box(slide1, "OPERAÇÃO LAVA JATO", top=Inches(2), size=Pt(48))
add_body(slide1, [
    "Principais Envolvidos e Fluxo Financeiro",
    "",
    "Ministério Público Federal — Força-Tarefa de Curitiba",
    f"Procurador responsável: {nome_procurador}",
    f"Juízo: 13ª Vara Federal de Curitiba — Dr. {nome_juiz}",
    "",
    "DOCUMENTO SIGILOSO — USO INTERNO",
], top=Inches(3.5), size=Pt(18))

# ============================================================
# SLIDE 2 — Organograma
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2)
add_title_box(slide2, "ORGANOGRAMA DOS ENVOLVIDOS")
add_body(slide2, [
    "NÚCLEO POLÍTICO:",
    f"  • Ex-Presidente: {nome_expresidente} (CPF: {cpf_expresidente})",
    f"  • Ex-Ministro da Casa Civil: {nome_exministro1}",
    f"  • Ex-Ministro da Fazenda: {nome_exministro2}",
    "",
    "NÚCLEO PETROBRAS:",
    f"  • Dir. de Abastecimento: {nome_dir_petrobras1}",
    f"  • Dir. Internacional: {nome_dir_petrobras2}",
    f"  • Dir. de Serviços: {nome_dir_petrobras3}",
    "",
    "NÚCLEO FINANCEIRO (DOLEIROS):",
    f"  • {nome_doleiro1} (CPF: {cpf_doleiro1}) — operador principal",
    f"  • {nome_doleiro2} — câmbio paralelo em Nova York",
    "",
    "NÚCLEO EMPREITEIRAS:",
    f"  • {nome_empreiteiro1} — Presidente da {empresa1}",
    f"  • {nome_empreiteiro2} — CEO da {empresa2}",
    f"  • {nome_empreiteiro3} — Dir. Financeiro da {empresa3}",
], top=Inches(1.6), size=Pt(14))

# ============================================================
# SLIDE 3 — Linha do tempo
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3)
add_title_box(slide3, "LINHA DO TEMPO — 2003 A 2016")
add_body(slide3, [
    "2003 — Início do esquema de propinas na Diretoria de Abastecimento",
    f"2004 — {nome_dir_petrobras1} assume diretoria; pagamentos via {nome_doleiro1}",
    f"2005 — {empresa1} (CNPJ: {cnpj1}) fecha contrato superfaturado de R$ 890 milhões",
    f"2006 — {nome_doleiro2} abre contas na Suíça para receber propinas",
    "2008 — Compra da refinaria de Pasadena com sobrepreço de R$ 1,18 bilhão",
    f"2009 — {nome_empreiteiro2} repassa R$ 75 milhões a partidos via caixa dois",
    f"2012 — {nome_dir_petrobras3} autoriza contratos sem licitação totalizando R$ 2,1 bilhões",
    "2014-03-17 — Deflagração da Fase 1 (Operação Lava Jato)",
    f"2014-11-14 — Prisão de {nome_doleiro1} na Fase 7",
    f"2015-02-06 — Delação premiada de {nome_delator1} (CPF: {cpf_delator1})",
    f"2015-06-19 — Prisão de {nome_empreiteiro1} na Fase 14",
    f"2015-08-20 — Condenação de {nome_dir_petrobras1}: 12 anos e 2 meses",
    f"2016-03-04 — Condução coercitiva de {nome_expresidente}",
    f"2016-09-12 — {nome_expresidente} condenado a 9 anos e 6 meses",
], top=Inches(1.6), size=Pt(13))

# ============================================================
# SLIDE 4 — Fluxo do dinheiro
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4)
add_title_box(slide4, "FLUXO DO DINHEIRO — VALORES EM R$")
add_body(slide4, [
    "ORIGEM → DESTINO → VALOR",
    "",
    f"{empresa1} → {nome_doleiro1} → R$ 187.500.000,00 (propina Abastecimento)",
    f"{empresa2} → {nome_doleiro2} → R$ 93.200.000,00 (propina Internacional)",
    f"{empresa3} → Caixa dois partidário → R$ 142.800.000,00",
    f"{empresa4} → {nome_expresidente} (via instituto) → R$ 3.700.000,00 (tríplex + sítio)",
    f"{nome_doleiro1} → Contas Suíça → R$ 256.400.000,00 (lavagem)",
    f"{nome_doleiro1} → {nome_dir_petrobras1} → R$ 25.600.000,00 (conta Mônaco)",
    f"{nome_doleiro2} → {nome_dir_petrobras2} → R$ 18.900.000,00 (Liechtenstein)",
    "",
    "TOTAL ESTIMADO DO DESVIO: R$ 6.194.000.000,00",
    "TOTAL RECUPERADO (até 2016): R$ 3.210.000.000,00",
], top=Inches(1.6), size=Pt(14))

# ============================================================
# SLIDE 5 — Empresas com CNPJs
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5)
add_title_box(slide5, "EMPRESAS ENVOLVIDAS")
add_body(slide5, [
    f"1. {empresa1}",
    f"   CNPJ: {cnpj1}",
    f"   Presidente: {nome_empreiteiro1}",
    f"   Endereço: {fake.address()}",
    f"   Contratos Petrobras: R$ 12,8 bilhões (2004-2014)",
    "",
    f"2. {empresa2}",
    f"   CNPJ: {cnpj2}",
    f"   CEO: {nome_empreiteiro2}",
    f"   Endereço: {fake.address()}",
    f"   Contratos Petrobras: R$ 8,3 bilhões (2005-2013)",
    "",
    f"3. {empresa3}",
    f"   CNPJ: {cnpj3}",
    f"   Dir. Financeiro: {nome_empreiteiro3}",
    f"   Endereço: {fake.address()}",
    f"   Contratos Petrobras: R$ 5,7 bilhões (2006-2014)",
    "",
    f"4. {empresa4}",
    f"   CNPJ: {cnpj4}",
    f"   Contratos Petrobras: R$ 3,1 bilhões (2008-2015)",
], top=Inches(1.6), size=Pt(12))

# ============================================================
# SLIDE 6 — Delatores
# ============================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide6)
add_title_box(slide6, "DELATORES E ACORDOS DE COLABORAÇÃO PREMIADA")
add_body(slide6, [
    f"DELATOR 1: {nome_delator1}",
    f"  CPF: {cpf_delator1}",
    f"  Cargo: Ex-diretor da {empresa1}",
    f"  Data do acordo: 06/02/2015",
    f"  Pena original: 15 anos → Pena reduzida: 3 anos (domiciliar)",
    f"  Revelou: esquema de propinas de R$ 187,5 milhões na Dir. Abastecimento",
    "",
    f"DELATOR 2: {nome_delator2}",
    f"  Cargo: Ex-gerente da Petrobras, indicado por {nome_exministro1}",
    f"  Data do acordo: 14/08/2015",
    f"  Pena original: 20 anos → Pena reduzida: 5 anos (semiaberto)",
    f"  Revelou: repasses de R$ 93,2 milhões via contas offshore",
    "",
    f"DELATOR 3: {nome_delator3}",
    f"  Cargo: Ex-presidente da {empresa3}",
    f"  Data do acordo: 22/11/2015",
    f"  Pena original: 18 anos → Pena reduzida: 4 anos + multa de R$ 50 milhões",
    f"  Revelou: nomes de 78 políticos recebedores de caixa dois",
], top=Inches(1.6), size=Pt(13))

# ============================================================
# SLIDE 7 — Valores recuperados
# ============================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide7)
add_title_box(slide7, "VALORES RECUPERADOS POR FASE DA OPERAÇÃO")
add_body(slide7, [
    "FASE    |  DATA        |  RECUPERADO         |  PRESOS",
    "--------|--------------|---------------------|--------",
    "Fase 1  |  17/03/2014  |  R$ 28.000.000      |  4",
    "Fase 7  |  14/11/2014  |  R$ 256.400.000     |  12",
    "Fase 14 |  19/06/2015  |  R$ 410.000.000     |  8",
    "Fase 23 |  21/01/2016  |  R$ 1.180.000.000   |  15",
    "Fase 24 |  04/03/2016  |  R$ 750.000.000     |  6",
    "Fase 34 |  22/09/2016  |  R$ 586.000.000     |  3",
    "",
    "TOTAL ACUMULADO: R$ 3.210.400.000,00",
    f"TOTAL DE DENÚNCIAS: 76 ações penais",
    f"TOTAL DE CONDENADOS: 174 pessoas",
    f"MAIOR PENA INDIVIDUAL: {nome_dir_petrobras1} — 12 anos, 2 meses e 20 dias",
], top=Inches(1.6), size=Pt(14))

# ============================================================
# SLIDE 8 — Conclusões
# ============================================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide8)
add_title_box(slide8, "CONCLUSÕES E PEDIDOS AO JUÍZO")
add_body(slide8, [
    "1. DEMONSTRAÇÃO DO ESQUEMA:",
    "   Ficou comprovado o esquema sistêmico de desvio de recursos da Petrobras",
    "   envolvendo agentes políticos, diretores e operadores financeiros.",
    "",
    "2. PEDIDOS:",
    f"   a) Condenação de {nome_expresidente} por corrupção passiva e lavagem (arts. 317 e 1º da Lei 9.613/98)",
    f"   b) Condenação de {nome_dir_petrobras1}, {nome_dir_petrobras2} e {nome_dir_petrobras3} por corrupção (art. 317 CP)",
    f"   c) Condenação de {nome_empreiteiro1}, {nome_empreiteiro2} e {nome_empreiteiro3} por corrupção ativa (art. 333 CP)",
    f"   d) Condenação de {nome_doleiro1} e {nome_doleiro2} por lavagem de dinheiro (Lei 9.613/98)",
    f"   e) Perda de bens: R$ 728.600.000,00 em contas e imóveis apreendidos",
    f"   f) Reparação de danos à Petrobras: R$ 6.194.000.000,00",
    "",
    "3. FUNDAMENTAÇÃO LEGAL:",
    "   Arts. 317, 333 do CP; Lei 9.613/98; Lei 12.850/13 (organizações criminosas);",
    "   Lei 8.429/92 (improbidade); Lei 12.846/13 (anticorrupção empresarial)",
], top=Inches(1.6), size=Pt(13))

# --- Salvar ---
out = "data/real/lavajato_apresentacao.pptx"
prs.save(out)
print(f"PPTX criado: {out} ({os.path.getsize(out) / 1024:.1f} KB)")
print(f"8 slides com dados fictícios gerados por Faker pt_BR")
