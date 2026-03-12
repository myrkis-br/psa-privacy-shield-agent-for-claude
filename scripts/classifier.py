"""
PSA - Privacy Shield Agent
Módulo: classifier.py (v6.0 — Risk Engine, Módulo 1)

Classifica documentos quanto ao risco LGPD:
  - Tipo e subtipo do documento
  - Número estimado de titulares
  - Presença de dados sensíveis (Art. 5°, II / Art. 11)
  - Score de risco 1-10 (ANPD Resolução nº 4/2023)
  - Classificação ANPD: leve / média / grave
  - Recomendação de cobertura de anonimização
  - Flag de consulta cloud (quando local não é suficiente)

Uso:
  from classifier import classify_document
  result = classify_document("data/real/arquivo.csv")
"""

import json
import re
import sys
from pathlib import Path
from typing import Any, Dict, List, Optional, Set

import pandas as pd

# ---------------------------------------------------------------------------
# Configuração
# ---------------------------------------------------------------------------

BASE_DIR = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(BASE_DIR / "scripts"))

# ---------------------------------------------------------------------------
# Categorias sensíveis (Art. 5°, II / Art. 11 LGPD)
# ---------------------------------------------------------------------------

# Dados sensíveis stricto sensu (Art. 11 LGPD) → +3 no risk score
SENSITIVE_HEALTH = {
    "saude", "saúde", "health", "cid", "diagnostico", "diagnóstico",
    "tratamento", "medicamento", "prescricao", "prescrição", "exame",
    "laudo", "atestado", "prontuario", "prontuário", "alergia",
    "patologia", "cirurgia", "internacao", "internação",
}
SENSITIVE_BIOMETRIC = {
    "biometria", "biometric", "facial", "iris", "íris",
    "retina", "impressao_digital", "impressão_digital",
    "reconhecimento_facial", "assinatura_digital",
}
SENSITIVE_MINOR = {
    "menor", "crianca", "criança", "adolescente", "infantil",
    "pediatr", "escola", "aluno", "estudante",
}
SENSITIVE_RACIAL = {
    "raca", "raça", "etnia", "cor_pele", "cor_raca",
    "indigena", "indígena", "quilombola", "autodeclaracao",
}
# Keywords curtas que exigem word boundary (evita "cor" em "cooperativo")
_RACIAL_WB = {"etnia"}
SENSITIVE_POLITICAL = {
    "partido", "filiacao_partidaria", "filiação", "sindicato",
    "sindical", "greve",
}
SENSITIVE_RELIGIOUS = {
    "religiao", "religião", "credo", "culto", "igreja",
}
SENSITIVE_SEXUAL = {
    "orientacao_sexual", "orientação_sexual", "genero", "gênero",
    "sexualidade",
}

ALL_SENSITIVE = (
    SENSITIVE_HEALTH | SENSITIVE_BIOMETRIC | SENSITIVE_MINOR |
    SENSITIVE_RACIAL | SENSITIVE_POLITICAL | SENSITIVE_RELIGIOUS |
    SENSITIVE_SEXUAL
)

# Dados financeiros → +2 no risk score
FINANCIAL_KEYWORDS = {
    "salario", "salário", "salary", "remuneracao", "remuneração",
    "vencimento", "bruto", "liquido", "líquido", "desconto",
    "gratificacao", "gratificação", "adicional", "abono", "auxilio",
    "auxílio", "subsidio", "subsídio", "conta_bancaria", "conta_corrente",
    "agencia", "agência", "num_banco", "cod_banco",
    "pix", "cartao", "cartão", "saldo", "debito", "débito",
    "credito", "crédito", "faturamento", "receita", "despesa",
    "plano_de_saude", "plano_saude", "convenio", "convênio",
    "seguro_saude", "seguro", "previdencia", "previdência",
}

# Dados cadastrais → +1 no risk score
CADASTRAL_KEYWORDS = {
    "cpf", "cnpj", "rg", "cnh", "pis", "pasep", "nis", "nit",
    "ctps", "sus", "titulo_eleitor", "passaporte", "matricula",
    "matrícula", "endereco", "endereço", "address", "cep",
    "telefone", "celular", "email", "nome", "name",
}

# Dados profissionais/funcionais → +1 no risk score
PROFESSIONAL_KEYWORDS = {
    "cargo", "departamento", "admissao", "admissão", "demissao",
    "demissão", "funcao", "função", "lotacao", "lotação",
    "vinculo", "vínculo", "jornada", "regime",
    "data_ingresso", "data_nomeacao", "data_nomeação",
    "orgao_exercicio", "org_lotacao",
}

# Keywords que devem ser buscadas com word boundary (regex \b)
# para evitar falsos positivos de substring
_WB_KEYWORDS = {"crm", "rqe", "cid", "voz"}

# ---------------------------------------------------------------------------
# Mapeamento extensão → tipo/subtipo
# ---------------------------------------------------------------------------

_EXT_MAP = {
    ".csv":  ("planilha", "CSV"),
    ".xlsx": ("planilha", "Excel"),
    ".xls":  ("planilha", "Excel legado"),
    ".docx": ("documento", "Word"),
    ".txt":  ("documento", "Texto"),
    ".rtf":  ("documento", "RTF"),
    ".odt":  ("documento", "ODT"),
    ".pdf":  ("documento", "PDF"),
    ".pptx": ("apresentacao", "PowerPoint"),
    ".eml":  ("email", "EML"),
    ".msg":  ("email", "MSG"),
    ".json": ("json", "JSON"),
    ".xml":  ("xml", "XML"),
    ".html": ("documento", "HTML"),
    ".htm":  ("documento", "HTML"),
    ".yaml": ("config", "YAML"),
    ".yml":  ("config", "YAML"),
    ".sql":  ("database", "SQL"),
    ".log":  ("log", "Log"),
    ".vcf":  ("contato", "vCard"),
    ".parquet": ("planilha", "Parquet"),
}

# ---------------------------------------------------------------------------
# Detecção de subtipo por conteúdo
# ---------------------------------------------------------------------------

_SUBTIPO_PATTERNS = {
    "laudo_medico": [
        r"laudo\s+m[eé]dico", r"atestado\s+m[eé]dico", r"CID[\s-]?\d{1,2}",
        r"diagn[oó]stico", r"prescri[çc][aã]o", r"exame\s+f[ií]sico",
        r"CRM[/\\]", r"anamnese",
    ],
    "folha_pagamento": [
        r"remunera[çc][aã]o", r"sal[aá]rio", r"vencimento", r"desconto",
        r"folha\s+(?:de\s+)?pagamento", r"bruto", r"l[ií]quido",
        r"IRRF", r"INSS", r"FGTS",
    ],
    "cadastro_funcionarios": [
        r"matr[ií]cula", r"cargo", r"departamento", r"admiss[aã]o",
        r"data\s+(?:de\s+)?admiss", r"fun[çc][aã]o", r"jornada",
    ],
    "contrato": [
        r"contrato", r"cl[aá]usula", r"contratante", r"contratad[oa]",
        r"objeto\s+do\s+contrato", r"vig[eê]ncia",
    ],
    "email_corporativo": [
        r"From:", r"To:", r"Subject:", r"De:", r"Para:", r"Assunto:",
        r"Enviado\s+(?:em|por):",
    ],
    "relatorio_financeiro": [
        r"balan[çc]o", r"demonstra[çc]", r"DRE", r"fluxo\s+de\s+caixa",
        r"ativo", r"passivo", r"patrim[oô]nio",
    ],
    "processo_judicial": [
        r"\d{7}-\d{2}\.\d{4}\.\d\.\d{2}\.\d{4}",
        r"tribunal", r"ju[ií]z", r"senten[çc]a", r"a[çc][aã]o\s+(?:civil|penal|trabalhista)",
    ],
}


def _detect_subtipo(text_sample: str) -> Optional[str]:
    """Detecta subtipo do documento por padrões no conteúdo."""
    text_lower = text_sample.lower()
    scores = {}
    for subtipo, patterns in _SUBTIPO_PATTERNS.items():
        count = 0
        for pat in patterns:
            if re.search(pat, text_sample, re.IGNORECASE):
                count += 1
        if count >= 2:  # mínimo 2 matches para classificar
            scores[subtipo] = count
    if scores:
        return max(scores, key=scores.get)
    return None


# ---------------------------------------------------------------------------
# Detecção de categorias sensíveis no conteúdo
# ---------------------------------------------------------------------------

def _kw_in_text(keyword: str, text: str) -> bool:
    """
    Verifica se keyword está no texto.
    Para keywords curtas (<=4 chars), usa word boundary para evitar falsos positivos.
    Ex: "cid" não deve casar com "decidir", "cor" não deve casar com "cooperativo".
    """
    if len(keyword) <= 4:
        return bool(re.search(r'\b' + re.escape(keyword) + r'\b', text))
    return keyword in text


def _detect_categories(text_sample: str, columns: Optional[List[str]] = None) -> Dict[str, bool]:
    """
    Detecta quais categorias de dados sensíveis estão presentes.
    Analisa tanto o texto do conteúdo quanto os nomes de colunas (se planilha).
    """
    all_text = text_sample.lower()
    # Colunas separadas por espaço para word boundary
    col_text = ""
    if columns:
        col_text = " " + " ".join(c.lower() for c in columns) + " "
        all_text += col_text

    categories = {
        "saude": False,
        "biometria": False,
        "menor": False,
        "racial": False,
        "politico": False,
        "religioso": False,
        "sexual": False,
        "financeiro": False,
        "cadastral": False,
        "profissional": False,
    }

    # Verifica cada categoria com word boundary para keywords curtas
    for kw in SENSITIVE_HEALTH:
        if _kw_in_text(kw, all_text):
            categories["saude"] = True
            break

    for kw in SENSITIVE_BIOMETRIC:
        if _kw_in_text(kw, all_text):
            categories["biometria"] = True
            break

    for kw in SENSITIVE_MINOR:
        if _kw_in_text(kw, all_text):
            categories["menor"] = True
            break

    for kw in SENSITIVE_RACIAL:
        if _kw_in_text(kw, all_text):
            categories["racial"] = True
            break

    for kw in SENSITIVE_POLITICAL:
        if _kw_in_text(kw, all_text):
            categories["politico"] = True
            break

    for kw in SENSITIVE_RELIGIOUS:
        if _kw_in_text(kw, all_text):
            categories["religioso"] = True
            break

    for kw in SENSITIVE_SEXUAL:
        if _kw_in_text(kw, all_text):
            categories["sexual"] = True
            break

    for kw in FINANCIAL_KEYWORDS:
        if _kw_in_text(kw, all_text):
            categories["financeiro"] = True
            break
    # Padrões multi-palavra para financeiro
    if not categories["financeiro"]:
        fin_patterns = [
            r"plano\s+de\s+sa[uú]de", r"conv[eê]nio\s+m[eé]dico",
            r"seguro\s+sa[uú]de", r"conta\s+(?:corrente|poupan[çc]a|banc[aá]ria)",
        ]
        for pat in fin_patterns:
            if re.search(pat, all_text):
                categories["financeiro"] = True
                break

    for kw in CADASTRAL_KEYWORDS:
        if _kw_in_text(kw, all_text):
            categories["cadastral"] = True
            break

    for kw in PROFESSIONAL_KEYWORDS:
        if _kw_in_text(kw, all_text):
            categories["profissional"] = True
            break

    return categories


# ---------------------------------------------------------------------------
# Estimativa de titulares
# ---------------------------------------------------------------------------

def _estimate_holders_spreadsheet(filepath: Path) -> int:
    """Estima número de titulares em planilha (= número de linhas de dados)."""
    suffix = filepath.suffix.lower()
    try:
        if suffix == ".csv":
            # Conta linhas sem carregar tudo, tenta diferentes encodings
            for enc in ("utf-8", "latin-1", "cp1252"):
                try:
                    with open(filepath, "r", encoding=enc) as f:
                        count = sum(1 for _ in f) - 1  # -1 para header
                    return max(count, 0)
                except UnicodeDecodeError:
                    continue
            return 0
        elif suffix in (".xlsx", ".xls"):
            df = pd.read_excel(filepath, dtype=str, nrows=0)
            # Para XLSX, usa openpyxl para contar sem carregar
            try:
                from openpyxl import load_workbook
                wb = load_workbook(filepath, read_only=True)
                ws = wb.active
                count = ws.max_row - 1 if ws.max_row else 0
                wb.close()
                return max(count, 0)
            except Exception:
                df_full = pd.read_excel(filepath, dtype=str)
                return len(df_full)
        elif suffix == ".parquet":
            df = pd.read_parquet(filepath, engine="pyarrow")
            return len(df)
    except Exception:
        return 0


def _estimate_holders_text(text: str) -> int:
    """
    Estima número de titulares em texto livre.
    Conta CPFs únicos, ou nomes próprios únicos como proxy.
    """
    # Conta CPFs
    cpfs = set(re.findall(r'\d{3}\.?\d{3}\.?\d{3}[-–]?\d{2}', text))
    if cpfs:
        return len(cpfs)
    # Fallback: assume 1 titular para documentos individuais (laudos, contratos)
    return 1


# ---------------------------------------------------------------------------
# Detecção de terceiros identificáveis
# ---------------------------------------------------------------------------

def _count_third_parties(text: str) -> int:
    """
    Conta terceiros identificáveis no texto (pessoas mencionadas além do titular).
    Detecta padrões como: Dr. Nome, CRM/XX NNNNN, nomes em Mixed Case com 2+ palavras.
    """
    if not text:
        return 0

    # Conta registros profissionais (CRM, CREFITO, CREA, OAB)
    prof_regs = re.findall(
        r'\b(?:CRM|CREFITO|CREA|OAB|CRO|CRP|COREN|CRF)[/\\-]?\s*[A-Z]{2}\s*[\d.:]+',
        text
    )

    # Conta nomes precedidos de títulos (Dr., Dra., Eng., etc.)
    titled_names = re.findall(
        r'\b(?:Dr\.?a?|Eng\.?a?|Prof\.?a?|Sr\.?a?)\s+'
        r'[A-ZÁÉÍÓÚÀÃÕÂÊÔÜÇ][a-záéíóúàãõâêôüç]+'
        r'(?:\s+(?:da|de|do|das|dos)\s+)?'
        r'(?:\s+[A-ZÁÉÍÓÚÀÃÕÂÊÔÜÇ][a-záéíóúàãõâêôüç]+)*',
        text
    )

    # Usa set de nomes únicos para evitar duplicatas
    unique = set()
    for name in titled_names:
        # Normaliza para comparação
        unique.add(name.strip().lower())
    for reg in prof_regs:
        unique.add(reg.strip().lower())

    return len(unique)


# ---------------------------------------------------------------------------
# JSON: extração de chaves, subtipo, titulares
# ---------------------------------------------------------------------------

def _load_json_safe(filepath: Path) -> Optional[Any]:
    """Carrega JSON com fallback de encoding."""
    for enc in ("utf-8", "latin-1", "cp1252"):
        try:
            with open(filepath, "r", encoding=enc) as f:
                return json.load(f)
        except (UnicodeDecodeError, json.JSONDecodeError):
            continue
    return None


def _collect_json_keys(obj: Any, prefix: str = "") -> Set[str]:
    """Coleta todas as chaves de um JSON recursivamente."""
    keys = set()  # type: Set[str]
    if isinstance(obj, dict):
        for k, v in obj.items():
            keys.add(k.lower())
            keys |= _collect_json_keys(v, k)
    elif isinstance(obj, list):
        for item in obj[:20]:  # amostra de 20 elementos
            keys |= _collect_json_keys(item, prefix)
    return keys


def _detect_json_subtipo(keys: Set[str]) -> str:
    """Detecta subtipo do JSON pelas chaves presentes."""
    # Adiciona fragmentos das chaves para matching
    expanded = set(keys)
    for k in keys:
        for part in re.split(r'[_\-]', k):
            if len(part) >= 2:
                expanded.add(part)
    keys = expanded

    scores = {}  # type: Dict[str, int]

    # Log de pagamentos / transações
    payment_keys = {"transaction", "payment", "card", "card_last4",
                    "pix_key", "chave_pix", "amount", "refund", "order_id",
                    "pix", "bank", "banco"}
    scores["log_pagamentos"] = len(keys & payment_keys)

    # Dados de saúde / pacientes
    health_keys = {"patient", "paciente", "health", "saude", "saúde",
                   "diagnostico", "cid", "prontuario", "prescricao"}
    scores["dados_saude"] = len(keys & health_keys)

    # Dados de RH / funcionários
    hr_keys = {"employee", "funcionario", "salary", "salario", "salário",
               "cargo", "departamento", "admissao", "rh", "folha"}
    scores["dados_rh"] = len(keys & hr_keys)

    # Dados pessoais genéricos
    personal_keys = {"user", "usuario", "cpf", "email", "name", "nome",
                     "phone", "telefone", "address", "endereco"}
    scores["dados_pessoais"] = len(keys & personal_keys)

    # Mínimo de 2 matches para classificar
    best = max(scores, key=scores.get)
    if scores[best] >= 2:
        return best
    return "json_generico"


def _estimate_holders_json(data: Any) -> int:
    """Estima titulares em JSON: len(array) ou 1 para objeto."""
    if isinstance(data, list):
        return len(data)
    return 1


def _extract_text_sample_json(data: Any, max_chars: int = 10000) -> str:
    """Extrai amostra de texto de um JSON para análise de categorias."""
    # Serializa uma amostra do JSON como texto
    if isinstance(data, list):
        sample = data[:20]
    else:
        sample = data
    text = json.dumps(sample, ensure_ascii=False, indent=0)
    return text[:max_chars]


# ---------------------------------------------------------------------------
# XML: extração de tags, subtipo, titulares
# ---------------------------------------------------------------------------

def _load_xml_safe(filepath: Path):
    """Carrega XML com ElementTree. Retorna (tree, root) ou (None, None)."""
    try:
        import xml.etree.ElementTree as ETree
        tree = ETree.parse(filepath)
        return tree, tree.getroot()
    except Exception:
        return None, None


def _strip_ns(tag: str) -> str:
    """Remove namespace de tag XML."""
    if "}" in tag:
        return tag.split("}", 1)[1]
    return tag


def _collect_xml_tags(element, tags: Optional[Set[str]] = None) -> Set[str]:
    """Coleta todos os nomes de tags de um XML recursivamente."""
    if tags is None:
        tags = set()
    tags.add(_strip_ns(element.tag).lower())
    for child in element:
        _collect_xml_tags(child, tags)
    return tags


def _detect_xml_subtipo(tags: Set[str]) -> str:
    """Detecta subtipo do XML pelas tags presentes."""
    # NF-e / NFC-e
    nfe_tags = {"nfeproc", "nfe", "infnfe", "emit", "dest", "det", "icmstot"}
    if len(tags & nfe_tags) >= 4:
        return "nfe"
    # CT-e
    cte_tags = {"cteproc", "cte", "infcte", "rem", "receb", "modal"}
    if len(tags & cte_tags) >= 3:
        return "cte"
    # NFS-e
    nfse_tags = {"nfse", "infnfse", "servico", "prestador", "tomador"}
    if len(tags & nfse_tags) >= 3:
        return "nfse"
    # SPED
    sped_tags = {"efd", "registro", "bloco", "contribuinte"}
    if len(tags & sped_tags) >= 2:
        return "sped"
    return "xml_generico"


def _estimate_holders_xml(root) -> int:
    """Estima titulares em XML fiscal."""
    count = 0
    for elem in root.iter():
        local = _strip_ns(elem.tag).lower()
        # Cada CPF/CNPJ em emit/dest/transp é um titular
        if local in ("cpf", "cnpj") and elem.text and elem.text.strip():
            count += 1
    return max(count, 1)


def _extract_text_sample_xml(root, max_chars: int = 10000) -> str:
    """Extrai amostra de texto de XML para análise de categorias."""
    parts = []
    for elem in root.iter():
        local = _strip_ns(elem.tag)
        if elem.text and elem.text.strip():
            parts.append(f"{local} {elem.text.strip()}")
        if len(" ".join(parts)) >= max_chars:
            break
    return " ".join(parts)[:max_chars]


# ---------------------------------------------------------------------------
# Extração de amostra de texto para análise
# ---------------------------------------------------------------------------

def _extract_text_sample(filepath: Path, max_chars: int = 10000) -> str:
    """Extrai amostra de texto do arquivo para classificação."""
    suffix = filepath.suffix.lower()

    if suffix == ".csv":
        try:
            for enc in ("utf-8-sig", "utf-8", "latin-1", "cp1252"):
                for sep in (",", ";", "\t", "|"):
                    try:
                        df = pd.read_csv(filepath, dtype=str, encoding=enc, sep=sep, nrows=20)
                        if len(df.columns) > 1:
                            text = " ".join(str(c) for c in df.columns) + "\n"
                            text += df.head(10).to_string()
                            return text[:max_chars]
                    except Exception:
                        continue
        except Exception:
            pass
        return ""

    elif suffix in (".xlsx", ".xls"):
        try:
            df = pd.read_excel(filepath, dtype=str, nrows=20)
            text = " ".join(str(c) for c in df.columns) + "\n"
            text += df.head(10).to_string()
            return text[:max_chars]
        except Exception:
            return ""

    elif suffix in (".txt", ".docx", ".rtf", ".odt"):
        try:
            if suffix == ".txt":
                for enc in ("utf-8", "latin-1", "cp1252"):
                    try:
                        with open(filepath, "r", encoding=enc) as f:
                            return f.read(max_chars)
                    except UnicodeDecodeError:
                        continue
            elif suffix == ".docx":
                from docx import Document
                doc = Document(str(filepath))
                text = "\n".join(p.text for p in doc.paragraphs[:50])
                return text[:max_chars]
            elif suffix == ".rtf":
                from striprtf.striprtf import rtf_to_text
                for enc in ("utf-8", "latin-1", "cp1252"):
                    try:
                        raw = filepath.read_text(encoding=enc)
                        return rtf_to_text(raw)[:max_chars]
                    except (UnicodeDecodeError, UnicodeError):
                        continue
            elif suffix == ".odt":
                from odf.opendocument import load as odf_load
                doc = odf_load(str(filepath))
                parts = []
                for elem in doc.text.childNodes:
                    if hasattr(elem, 'qname'):
                        tag = elem.qname[1] if isinstance(elem.qname, tuple) else str(elem.qname)
                        if tag in ('p', 'h'):
                            # Extract text recursively
                            def _get_text(el):
                                result = []
                                if hasattr(el, 'childNodes'):
                                    for ch in el.childNodes:
                                        if hasattr(ch, 'data'):
                                            result.append(ch.data)
                                        else:
                                            result.append(_get_text(ch))
                                return ''.join(result)
                            text = _get_text(elem).strip()
                            if text:
                                parts.append(text)
                return "\n".join(parts)[:max_chars]
        except Exception:
            return ""

    elif suffix == ".pdf":
        try:
            import pdfplumber
            text = ""
            with pdfplumber.open(filepath) as pdf:
                for page in pdf.pages[:5]:
                    text += (page.extract_text() or "") + "\n"
                    if len(text) >= max_chars:
                        break
            return text[:max_chars]
        except Exception:
            return ""

    elif suffix in (".eml", ".msg"):
        try:
            if suffix == ".eml":
                import email
                with open(filepath, "r", encoding="utf-8", errors="replace") as f:
                    msg = email.message_from_file(f)
                parts = []
                if msg.is_multipart():
                    for part in msg.walk():
                        if part.get_content_type() == "text/plain":
                            payload = part.get_payload(decode=True)
                            if payload:
                                parts.append(payload.decode("utf-8", errors="replace"))
                else:
                    payload = msg.get_payload(decode=True)
                    if payload:
                        parts.append(payload.decode("utf-8", errors="replace"))
                # Include headers
                header_text = f"From: {msg.get('From', '')} To: {msg.get('To', '')} Subject: {msg.get('Subject', '')}"
                return (header_text + "\n" + "\n".join(parts))[:max_chars]
            elif suffix == ".msg":
                import extract_msg
                msg = extract_msg.Message(str(filepath))
                text = f"From: {msg.sender} To: {msg.to} Subject: {msg.subject}\n{msg.body}"
                msg.close()
                return text[:max_chars]
        except Exception:
            return ""

    elif suffix == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(str(filepath))
            text = ""
            for slide in prs.slides[:10]:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text[:max_chars]
        except Exception:
            return ""

    elif suffix == ".json":
        data = _load_json_safe(filepath)
        if data is not None:
            return _extract_text_sample_json(data, max_chars)
        return ""

    elif suffix == ".xml":
        _, root = _load_xml_safe(filepath)
        if root is not None:
            return _extract_text_sample_xml(root, max_chars)
        return ""

    elif suffix in (".html", ".htm"):
        try:
            import re as _re
            for enc in ("utf-8", "latin-1", "cp1252"):
                try:
                    raw = filepath.read_text(encoding=enc)
                    # Strip HTML tags to get text
                    text = _re.sub(r'<(script|style)[^>]*>.*?</\1>', ' ', raw, flags=_re.DOTALL | _re.IGNORECASE)
                    text = _re.sub(r'<[^>]+>', ' ', text)
                    text = _re.sub(r'\s+', ' ', text).strip()
                    return text[:max_chars]
                except UnicodeDecodeError:
                    continue
        except Exception:
            return ""

    elif suffix in (".yaml", ".yml"):
        try:
            for enc in ("utf-8", "latin-1"):
                try:
                    return filepath.read_text(encoding=enc)[:max_chars]
                except UnicodeDecodeError:
                    continue
        except Exception:
            return ""

    elif suffix == ".sql":
        try:
            for enc in ("utf-8", "latin-1", "cp1252"):
                try:
                    return filepath.read_text(encoding=enc)[:max_chars]
                except UnicodeDecodeError:
                    continue
        except Exception:
            return ""

    elif suffix == ".log":
        try:
            for enc in ("utf-8", "latin-1", "cp1252"):
                try:
                    return filepath.read_text(encoding=enc)[:max_chars]
                except UnicodeDecodeError:
                    continue
        except Exception:
            return ""

    elif suffix == ".vcf":
        try:
            return filepath.read_text(encoding="utf-8", errors="replace")[:max_chars]
        except Exception:
            return ""

    elif suffix == ".parquet":
        try:
            df = pd.read_parquet(filepath, engine="pyarrow")
            text = " ".join(str(c) for c in df.columns) + "\n"
            text += df.head(20).to_string()
            return text[:max_chars]
        except Exception:
            return ""

    return ""


def _get_columns(filepath: Path) -> Optional[List[str]]:
    """Retorna nomes de colunas se for planilha, None caso contrário."""
    suffix = filepath.suffix.lower()
    if suffix == ".csv":
        for enc in ("utf-8-sig", "utf-8", "latin-1", "cp1252"):
            for sep in (",", ";", "\t", "|"):
                try:
                    df = pd.read_csv(filepath, dtype=str, encoding=enc, sep=sep, nrows=0)
                    if len(df.columns) > 1:
                        return list(df.columns)
                except Exception:
                    continue
    elif suffix in (".xlsx", ".xls"):
        try:
            df = pd.read_excel(filepath, dtype=str, nrows=0)
            return list(df.columns)
        except Exception:
            pass
    elif suffix == ".parquet":
        try:
            df = pd.read_parquet(filepath, engine="pyarrow")
            return list(df.columns)
        except Exception:
            pass
    return None


# ---------------------------------------------------------------------------
# Função principal: classify_document
# ---------------------------------------------------------------------------

def classify_document(filepath) -> Dict:
    """
    Classifica um documento quanto ao risco LGPD.

    Args:
        filepath: caminho para o arquivo (str ou Path)

    Returns:
        Dict com:
            tipo: str — tipo do documento (planilha, documento, email, etc.)
            subtipo: str — subtipo detectado (laudo_medico, folha_pagamento, etc.)
            n_titulares_estimado: int — número estimado de titulares de dados
            tem_sensivel: bool — se contém dados sensíveis (Art. 11 LGPD)
            categorias_sensiveis: List[str] — categorias detectadas
            risk_score: int — score de risco 1-10
            classificacao_anpd: str — leve / média / grave
            cobertura_recomendada: str — descrição da cobertura recomendada
            consulta_cloud: bool — se deve consultar API cloud
            justificativa: str — explicação do score
    """
    filepath = Path(filepath)
    suffix = filepath.suffix.lower()

    # --- Tipo e subtipo ---
    tipo, subtipo_label = _EXT_MAP.get(suffix, ("desconhecido", "desconhecido"))

    # --- Extração de texto e colunas ---
    text_sample = _extract_text_sample(filepath)
    columns = _get_columns(filepath)

    # --- JSON: extração de chaves + subtipo por estrutura ---
    json_keys = set()  # type: Set[str]
    json_data = None
    if suffix == ".json":
        json_data = _load_json_safe(filepath)
        if json_data is not None:
            json_keys = _collect_json_keys(json_data)

    # --- XML: extração de tags + subtipo por estrutura ---
    xml_tags = set()  # type: Set[str]
    xml_root = None
    if suffix == ".xml":
        _, xml_root = _load_xml_safe(filepath)
        if xml_root is not None:
            xml_tags = _collect_xml_tags(xml_root)

    # --- Subtipo por conteúdo ---
    if suffix == ".json" and json_keys:
        subtipo = _detect_json_subtipo(json_keys)
    elif suffix == ".xml" and xml_tags:
        subtipo = _detect_xml_subtipo(xml_tags)
    else:
        subtipo = _detect_subtipo(text_sample)
        if subtipo is None:
            subtipo = subtipo_label.lower().replace(" ", "_")

    # --- Categorias sensíveis ---
    # Para JSON/XML, passa tags/chaves e seus fragmentos como "colunas"
    effective_columns = columns
    if suffix == ".json" and json_keys:
        fragments = set()  # type: Set[str]
        for k in json_keys:
            fragments.add(k)
            for part in re.split(r'[_\-]', k):
                if len(part) >= 2:
                    fragments.add(part)
        effective_columns = list(fragments)
    elif suffix == ".xml" and xml_tags:
        effective_columns = list(xml_tags)
    categories = _detect_categories(text_sample, effective_columns)
    categorias_ativas = [k for k, v in categories.items() if v]

    # Dados sensíveis stricto sensu (Art. 11)
    art11_cats = {"saude", "biometria", "menor", "racial", "politico", "religioso", "sexual"}
    tem_sensivel = any(categories.get(c, False) for c in art11_cats)

    # --- Estimativa de titulares ---
    if suffix == ".json" and json_data is not None:
        n_titulares = _estimate_holders_json(json_data)
    elif suffix == ".xml" and xml_root is not None:
        n_titulares = _estimate_holders_xml(xml_root)
    elif tipo == "planilha":
        n_titulares = _estimate_holders_spreadsheet(filepath)
    else:
        n_titulares = _estimate_holders_text(text_sample)

    # --- Detecção de terceiros (pessoas mencionadas além do titular) ---
    n_terceiros = _count_third_parties(text_sample)

    # --- Cálculo do risk_score ---
    risk_score = 0
    justificativas = []

    # dado sensível (saúde, menor, biometria): +3
    if tem_sensivel:
        risk_score += 3
        cats_label = [c for c in art11_cats if categories.get(c)]
        justificativas.append(f"dado sensível Art. 11 ({', '.join(cats_label)}): +3")

    # dado financeiro: +2
    if categories.get("financeiro"):
        risk_score += 2
        justificativas.append("dado financeiro (salário/conta/PIX): +2")

    # dado cadastral: +1
    if categories.get("cadastral"):
        risk_score += 1
        justificativas.append("dado cadastral (CPF/RG/endereço): +1")

    # dado profissional/funcional: +1
    if categories.get("profissional"):
        risk_score += 1
        justificativas.append("dado profissional (cargo/função/lotação): +1")

    # volume de titulares
    if n_titulares <= 50:
        vol_bonus = 0
    elif n_titulares <= 500:
        vol_bonus = 1
    elif n_titulares <= 10000:
        vol_bonus = 2
    else:
        vol_bonus = 3
    if vol_bonus > 0:
        risk_score += vol_bonus
        justificativas.append(f"volume ({n_titulares} titulares): +{vol_bonus}")

    # terceiros identificáveis: +1 se >= 3 terceiros
    if n_terceiros >= 3:
        risk_score += 1
        justificativas.append(f"terceiros identificáveis ({n_terceiros} pessoas): +1")

    # cadastro funcional em volume (servidor público, RH): +1
    if (subtipo in ("cadastro_funcionarios", "folha_pagamento")
            and n_titulares > 100):
        risk_score += 1
        justificativas.append(f"cadastro funcional em volume ({subtipo}): +1")

    # laudo médico completo (diagnóstico + prescrição + atestado): +1
    if subtipo == "laudo_medico":
        med_indicators = 0
        text_lower = text_sample.lower()
        if re.search(r"cid[\s-]?\d{1,2}", text_lower):
            med_indicators += 1
        if re.search(r"prescri[çc][aã]o|medicament", text_lower):
            med_indicators += 1
        if re.search(r"atestado|afastamento", text_lower):
            med_indicators += 1
        if med_indicators >= 2:
            risk_score += 1
            justificativas.append(f"laudo médico completo (CID + prescrição/atestado): +1")

    # XML fiscal (NF-e/CT-e): dados de múltiplas partes → +2
    if suffix == ".xml" and xml_tags:
        fiscal_subtypes = {"nfe", "nfse", "cte"}
        if subtipo in fiscal_subtypes:
            # NF-e tem emitente + destinatário + transportadora = 3 partes
            parties = {"emit", "dest", "transporta", "rem", "receb",
                       "tomador", "prestador"}
            n_parties = len(xml_tags & parties)
            if n_parties >= 2:
                risk_score += 2
                justificativas.append(
                    f"documento fiscal ({subtipo}) com {n_parties} partes identificadas: +2"
                )
            # NF-e é inerentemente financeira (cobrança, fatura, impostos)
            if not categories.get("financeiro"):
                categories["financeiro"] = True
                if "financeiro" not in categorias_ativas:
                    categorias_ativas.append("financeiro")
                risk_score += 2
                justificativas.append("documento fiscal com valores/impostos: +2")
        # XML com endereços completos (xLgr + xBairro + xMun + CEP) → +1
        addr_tags = {"xlgr", "xbairro", "xmun", "cep"}
        if len(xml_tags & addr_tags) >= 3:
            risk_score += 1
            justificativas.append("endereços completos (logradouro+bairro+município+CEP): +1")

    # JSON com dados de rastreamento digital (IP, fingerprint, device): +1
    if suffix == ".json" and json_keys:
        tracking_keys = {"ip", "ip_address", "fingerprint", "device",
                         "user_agent", "session_id"}
        if json_keys & tracking_keys:
            risk_score += 1
            justificativas.append("rastreamento digital (IP/fingerprint/device): +1")

    # JSON: dados pessoais combinados — risco multiplicativo
    # Quanto mais tipos de PII coexistem, maior o risco de reidentificação
    if suffix == ".json" and json_keys:
        # Conta categorias de PII distintas nas chaves
        pii_groups = {
            "identificador": {"cpf", "cnpj", "rg", "cnh", "pis"},
            "contato": {"email", "phone", "telefone", "celular"},
            "localizacao": {"address", "endereco", "cep", "ip"},
            "financeiro": {"pix", "card", "cartao", "bank", "amount"},
            "rastreamento": {"fingerprint", "device", "user_agent"},
        }
        # Fragmenta chaves para matching
        key_fragments = set()
        for k in json_keys:
            key_fragments.add(k)
            for part in re.split(r'[_\-]', k):
                if len(part) >= 2:
                    key_fragments.add(part)
        pii_count = sum(
            1 for group_keys in pii_groups.values()
            if key_fragments & group_keys
        )
        if pii_count >= 4:
            bonus = pii_count - 2  # 4 tipos → +2, 5 tipos → +3
            risk_score += bonus
            justificativas.append(
                f"dados pessoais combinados ({pii_count} categorias PII): +{bonus}"
            )

    # formato desconhecido: +1
    if tipo == "desconhecido":
        risk_score += 1
        justificativas.append("formato desconhecido: +1")

    # Garantir mínimo 1 se há qualquer dado pessoal
    if risk_score == 0 and (categories.get("cadastral") or categories.get("financeiro")):
        risk_score = 1
        justificativas.append("mínimo: dado pessoal detectado")

    # Clamp 1-10
    risk_score = max(1, min(10, risk_score))

    # --- Classificação ANPD (Resolução nº 4/2023) ---
    if risk_score <= 3:
        classificacao = "leve"
    elif risk_score <= 6:
        classificacao = "média"
    else:
        classificacao = "grave"

    # --- Cobertura recomendada ---
    if classificacao == "grave":
        cobertura = "anonimização completa + RIPD obrigatório + encarregado deve ser notificado"
    elif classificacao == "média":
        cobertura = "anonimização completa + RIPD recomendado"
    else:
        cobertura = "anonimização padrão"

    # --- Consulta cloud ---
    # consulta_cloud = True quando risk_score >= 4 AND taxa_sucesso_local < 95%
    # Como não temos taxa_sucesso_local neste ponto, usamos heurística:
    # Se o subtipo é desconhecido ou se tem dados sensíveis com score alto, sugere cloud
    consulta_cloud = (
        risk_score >= 4
        and (subtipo in (subtipo_label.lower().replace(" ", "_"),) or tem_sensivel)
    )

    justificativa = " | ".join(justificativas) if justificativas else "nenhum dado pessoal detectado"

    return {
        "tipo": tipo,
        "subtipo": subtipo,
        "n_titulares_estimado": n_titulares,
        "tem_sensivel": tem_sensivel,
        "categorias_sensiveis": categorias_ativas,
        "risk_score": risk_score,
        "classificacao_anpd": classificacao,
        "cobertura_recomendada": cobertura,
        "consulta_cloud": consulta_cloud,
        "justificativa": justificativa,
    }


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main():
    import json

    if len(sys.argv) < 2:
        print("Uso: python3 scripts/classifier.py <arquivo>")
        sys.exit(1)

    filepath = Path(sys.argv[1])
    if not filepath.exists():
        print(f"Arquivo não encontrado: {filepath}")
        sys.exit(1)

    result = classify_document(filepath)

    print("\n" + "=" * 70)
    print("PSA — CLASSIFICAÇÃO DE RISCO LGPD")
    print("=" * 70)
    print(f"  Tipo              : {result['tipo']}")
    print(f"  Subtipo           : {result['subtipo']}")
    print(f"  Titulares estimados: {result['n_titulares_estimado']}")
    print(f"  Dado sensível     : {'SIM' if result['tem_sensivel'] else 'NÃO'}")
    print(f"  Categorias        : {', '.join(result['categorias_sensiveis']) or 'nenhuma'}")
    print(f"  Risk Score        : {result['risk_score']}/10")
    print(f"  Classificação ANPD: {result['classificacao_anpd'].upper()}")
    print(f"  Cobertura         : {result['cobertura_recomendada']}")
    print(f"  Consulta Cloud    : {'SIM' if result['consulta_cloud'] else 'NÃO'}")
    print(f"  Justificativa     : {result['justificativa']}")
    print("=" * 70)

    # Salva JSON
    print(f"\n  JSON:\n{json.dumps(result, ensure_ascii=False, indent=2)}")


if __name__ == "__main__":
    main()
